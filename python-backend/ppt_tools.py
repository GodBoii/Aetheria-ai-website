import json
import logging
import mimetypes
import os
import re
import shutil
import subprocess
import tempfile
import uuid
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

from agno.agent import Agent
from agno.tools import Toolkit

from sandbox_persistence import get_persistence_service

logger = logging.getLogger(__name__)

PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

TEMPLATE_LAYOUTS = [
    {
        "type": "title",
        "name": "Cover",
        "usage": "Opening slide with large editorial title, subtitle, visual system panel, and optional metrics.",
    },
    {
        "type": "content",
        "name": "Insight Cards",
        "usage": "Claim-style slide with 3-4 short insight cards, side visual, and optional callout.",
    },
    {
        "type": "two_column",
        "name": "Comparison",
        "usage": "Two filled comparison panels; include left/right titles and bullets whenever possible.",
    },
    {
        "type": "chart",
        "name": "Evidence Chart",
        "usage": "Content slide with chart.data values for simple bar evidence.",
    },
    {
        "type": "table",
        "name": "Structured Table",
        "usage": "Content slide with table rows for comparison matrices, plans, or structured facts.",
    },
    {
        "type": "diagram",
        "name": "Process Flow",
        "usage": "Content slide with nodes or steps for workflows, systems, and timelines.",
    },
    {
        "type": "image",
        "name": "Visual Explanation",
        "usage": "Visual slide with image_path when available, otherwise a designed abstract visual and supporting cards.",
    },
]

TEMPLATES: Dict[str, Dict[str, Any]] = {
    "venture_blueprint": {
        "name": "Venture Blueprint",
        "description": "Premium pitch and business deck with bold left-rail titles, editorial image zones, and investor-grade evidence layouts.",
        "best_for": "pitch decks, business plans, strategy narratives, product launches",
        "design_brief": (
            "Use this as a true pitch/business storytelling system, not a recolored standard deck. "
            "Prefer strong claim-style titles, large left-side thesis blocks, right-side product or market visuals, "
            "metric chips, problem/solution contrast, market evidence charts, business model tables, and roadmap steps. "
            "Keep text concise so it fits: short titles, 3-4 bullets, and compact labels. Include image_path on visual/product slides when available."
        ),
        "layouts": TEMPLATE_LAYOUTS,
    },
    "aetheria_modern": {
        "name": "Aetheria Modern",
        "description": "Clean editorial deck for AI strategy and product narratives.",
        "best_for": "AI strategy, product plans, operational reviews",
        "design_brief": "Editorial, spacious, sharp blue/red/green accents, visual systems, insight cards, and restrained business polish.",
        "layouts": TEMPLATE_LAYOUTS,
    },
    "executive": {
        "name": "Executive Boardroom",
        "description": "Refined boardroom aesthetic with crisp data hierarchy.",
        "best_for": "business reviews, leadership updates, investor summaries",
        "design_brief": "Boardroom-ready, calm, metric-led, structured comparisons, polished evidence, and premium whitespace.",
        "layouts": TEMPLATE_LAYOUTS,
    },
    "startup_pitch": {
        "name": "Startup Pitch",
        "description": "High-contrast dark deck with bold metrics for investors.",
        "best_for": "startup fundraising, product launches, market narratives",
        "design_brief": "High-contrast, energetic, metric-first, product/storytelling slides with strong visual hierarchy.",
        "layouts": TEMPLATE_LAYOUTS,
    },
    "academic": {
        "name": "Academic Research",
        "description": "Formal scholarly layout with readable evidence and citations.",
        "best_for": "research talks, coursework, technical explainers",
        "design_brief": "Scholarly, readable, evidence-heavy, citation-friendly, with structured findings and methodology diagrams.",
        "layouts": TEMPLATE_LAYOUTS,
    },
    "creative_portfolio": {
        "name": "Creative Portfolio",
        "description": "Bold expressive deck with vibrant gradients and asymmetric layouts.",
        "best_for": "design portfolios, creative briefs, brand pitches",
        "design_brief": "Expressive, asymmetric, colorful, showcase-oriented, with strong visual slides and bold story beats.",
        "layouts": TEMPLATE_LAYOUTS,
    },
    "minimal_zen": {
        "name": "Minimal Zen",
        "description": "Ultra-clean whitespace design with restrained single-accent palette.",
        "best_for": "thought leadership, keynotes, minimalist reports",
        "design_brief": "Minimal, quiet, highly legible, generous whitespace, fewer words, and carefully paced emphasis.",
        "layouts": TEMPLATE_LAYOUTS,
    },
    "tech_dark": {
        "name": "Tech Neon",
        "description": "Dark engineering theme with electric neon accents and sharp edges.",
        "best_for": "technical demos, developer talks, product launches",
        "design_brief": "Dark technical interface style, neon accents, architecture/process diagrams, specs, and benchmark evidence.",
        "layouts": TEMPLATE_LAYOUTS,
    },
    "corporate_gradient": {
        "name": "Corporate Horizon",
        "description": "Professional gradient-rich deck with structured visual hierarchy.",
        "best_for": "quarterly reports, all-hands meetings, client proposals",
        "design_brief": "Corporate, confident, structured, client-ready, with KPI evidence, process slides, and organized summaries.",
        "layouts": TEMPLATE_LAYOUTS,
    },
}


def _template_summary(template_id: str, template: Dict[str, Any]) -> Dict[str, str]:
    return {
        "id": template_id,
        "name": str(template.get("name", template_id)),
        "best_for": str(template.get("best_for", "")),
        "description": str(template.get("description", "")),
    }


def _resolve_template_id(value: Any) -> Optional[str]:
    text = str(value or "").strip()
    if not text:
        return None
    normalized = re.sub(r"[^a-z0-9]+", "_", text.lower()).strip("_")
    if normalized in TEMPLATES:
        return normalized
    for template_id, template in TEMPLATES.items():
        candidates = {
            template_id,
            re.sub(r"[^a-z0-9]+", "_", str(template.get("name", "")).lower()).strip("_"),
        }
        if normalized in candidates:
            return template_id
    return None


def _safe_slug(value: str, fallback: str = "presentation") -> str:
    text = re.sub(r"[^a-zA-Z0-9._-]+", "-", str(value or "").strip().lower())
    text = re.sub(r"-+", "-", text).strip("-._")
    return (text or fallback)[:80]


def _parse_jsonish(value: Any, fallback: Any) -> Any:
    if value is None:
        return fallback
    if isinstance(value, (dict, list)):
        return value
    if isinstance(value, str):
        text = value.strip()
        if not text:
            return fallback
        try:
            return json.loads(text)
        except Exception:
            return fallback
    return fallback


def _normalize_slide(slide: Any, index: int, topic: str) -> Dict[str, Any]:
    if isinstance(slide, str):
        return {"type": "content", "title": f"Slide {index + 1}", "content": slide}
    if not isinstance(slide, dict):
        return {"type": "content", "title": f"Slide {index + 1}", "content": str(slide)}

    normalized = dict(slide)
    normalized.setdefault("type", "title" if index == 0 else "content")
    normalized.setdefault("title", topic if index == 0 else f"Slide {index + 1}")
    return normalized


class PresentationTools(Toolkit):
    """
    Native PowerPoint generator for Aetheria's presentation sub-agent.
    Uses pptxgenjs to create editable .pptx files; no HTML/CSS is used for the
    actual PowerPoint output.
    """

    def __init__(
        self,
        *,
        user_id: Optional[str],
        session_id: Optional[str],
        message_id: Optional[str],
        socketio=None,
        sid: Optional[str] = None,
    ):
        super().__init__(
            name="presentation_tools",
            tools=[
                self.create_presentation,
                self.list_presentation_templates,
                self.get_presentation_template_details,
                self.edit_presentation_text,
            ],
        )
        self.user_id = user_id
        self.session_id = session_id
        self.message_id = message_id
        self.socketio = socketio
        self.sid = sid
        self.backend_dir = Path(__file__).resolve().parent
        self.repo_root = self.backend_dir.parent
        self.renderer_path = self._resolve_renderer_path()

    def _resolve_renderer_path(self) -> Path:
        configured = os.getenv("PPTX_RENDERER_PATH")
        candidates = [
            Path(configured) if configured else None,
            self.backend_dir / "pptx-renderer.js",
            self.backend_dir / "js" / "pptx-renderer.js",
            self.repo_root / "js" / "pptx-renderer.js",
        ]
        for candidate in candidates:
            if candidate and candidate.exists():
                return candidate
        return self.backend_dir / "pptx-renderer.js"

    def list_presentation_templates(self) -> Dict[str, Any]:
        """List available native PowerPoint templates using compact summaries."""
        templates = [
            _template_summary(key, value)
            for key, value in TEMPLATES.items()
        ]
        return {
            "ok": True,
            "message": (
                "Available native PowerPoint templates. This is a compact list; "
                "call get_presentation_template_details(template_id) for layout and design details."
            ),
            "data": {"templates": templates},
            "metadata": {
                "kind": "presentation_tool_output",
                "action": "list_templates",
                "preview_type": "presentation_templates",
                "title": "Presentation templates",
                "inline": {
                    "templates": templates
                },
            },
        }

    def get_presentation_template_details(self, template_id: str) -> Dict[str, Any]:
        """Return detailed layout and design guidance for one presentation template."""
        resolved_id = _resolve_template_id(template_id)
        if not resolved_id:
            return self._error(
                f"Unknown presentation template '{template_id}'. "
                f"Available template ids: {', '.join(TEMPLATES.keys())}"
            )

        template = TEMPLATES[resolved_id]
        return {
            "ok": True,
            "message": f"Template details for {template.get('name', resolved_id)}.",
            "data": {
                "id": resolved_id,
                **template,
            },
            "metadata": {
                "kind": "presentation_tool_output",
                "action": "template_details",
                "preview_type": "text",
                "title": f"Template details: {template.get('name', resolved_id)}",
                "inline": {
                    "id": resolved_id,
                    "name": template.get("name", resolved_id),
                    "best_for": template.get("best_for", ""),
                    "description": template.get("description", ""),
                },
            },
        }

    def create_presentation(
        self,
        topic: str,
        slides: Any,
        template: str = "aetheria_modern",
        filename: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Create a native editable .pptx from structured slide definitions.

        Args:
            topic: Presentation title or topic.
            slides: List of slide dictionaries, or a JSON string. Supported slide
                types include title/cover, content, two_column/comparison, chart,
                table, diagram/process, image/visual. Slides should include
                structured fields whenever possible: bullets/content, left/right
                comparison content, metrics, chart.data, table rows, nodes/steps,
                callout, notes, captions, visual_summary, or image_path.
            template: One of venture_blueprint, aetheria_modern, executive, startup_pitch, academic, creative_portfolio, minimal_zen, tech_dark, corporate_gradient.
            filename: Optional output filename ending in .pptx.
        """
        try:
            if not topic or not str(topic).strip():
                return self._error("topic is required")

            slide_list = _parse_jsonish(slides, [])
            if not isinstance(slide_list, list):
                return self._error("slides must be a list of slide dictionaries or a JSON list")
            if not slide_list:
                slide_list = [
                    {"type": "title", "title": topic, "subtitle": "Generated by Aetheria AI"},
                    {"type": "content", "title": "Key points", "bullets": ["Main idea", "Supporting proof", "Next step"]},
                ]

            template_id = _resolve_template_id(template)
            if not template_id:
                return self._error(
                    f"Unknown presentation template '{template}'. "
                    f"Available template ids: {', '.join(TEMPLATES.keys())}"
                )
            safe_name = _safe_slug(filename or topic)
            if not safe_name.endswith(".pptx"):
                safe_name = f"{safe_name}.pptx"

            work_dir = Path(tempfile.mkdtemp(prefix="aetheria-ppt-"))
            output_path = work_dir / safe_name
            payload_path = work_dir / "payload.json"
            normalized_slides = [
                _normalize_slide(slide, index, str(topic))
                for index, slide in enumerate(slide_list)
            ]

            payload = {
                "topic": str(topic).strip(),
                "slides": normalized_slides,
                "template": template_id,
                "output_path": str(output_path),
            }
            payload_path.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")

            node = shutil.which("node") or shutil.which("node.exe")
            if not node:
                return self._error("Node.js was not found; presentation rendering requires Node.js and pptxgenjs.")
            if not self.renderer_path.exists():
                return self._error(f"PowerPoint renderer not found at {self.renderer_path}")

            completed = subprocess.run(
                [node, str(self.renderer_path), str(payload_path)],
                cwd=str(self.backend_dir),
                capture_output=True,
                text=True,
                timeout=120,
            )
            renderer_stdout = (completed.stdout or "").strip()
            renderer_result = _parse_jsonish(renderer_stdout, {})
            if completed.returncode != 0 or not renderer_result.get("ok"):
                logger.error(
                    "PPT renderer failed rc=%s stdout=%s stderr=%s",
                    completed.returncode,
                    renderer_stdout[:2000],
                    (completed.stderr or "")[:2000],
                )
                return self._error(
                    renderer_result.get("error")
                    or (completed.stderr or "PowerPoint renderer failed").strip()
                )

            pptx_bytes = output_path.read_bytes()
            artifact_id = self._persist_pptx(
                file_path=str(output_path),
                filename=safe_name,
                file_content=pptx_bytes,
            )
            download_url = None
            if artifact_id and self.user_id:
                try:
                    download_url = get_persistence_service().get_artifact_download_url(
                        artifact_id=artifact_id,
                        user_id=str(self.user_id),
                        expiry=3600,
                    )
                except Exception as exc:
                    logger.warning("Unable to generate presentation download URL: %s", exc)

            metadata = {
                "kind": "presentation_tool_output",
                "action": "create_presentation",
                "preview_type": "presentation",
                "output_id": artifact_id or str(uuid.uuid4()),
                "artifact_id": artifact_id,
                "title": str(topic).strip(),
                "summary": f"Created {len(normalized_slides)}-slide native PowerPoint deck.",
                "filename": safe_name,
                "mime_type": PPTX_MIME_TYPE,
                "download_url": download_url,
                "template": renderer_result.get("template"),
                "layout_validation": renderer_result.get("layout_validation"),
                "inline": {
                    "topic": str(topic).strip(),
                    "slide_count": len(normalized_slides),
                    "size_bytes": len(pptx_bytes),
                    "slides": renderer_result.get("slides") or [],
                },
            }
            self._emit_presentation_created(metadata)
            return {
                "ok": True,
                "message": (
                    f"Created native editable PowerPoint '{safe_name}' "
                    f"with {len(normalized_slides)} slides."
                ),
                "data": {
                    "artifact_id": artifact_id,
                    "filename": safe_name,
                    "download_url": download_url,
                    "local_path": str(output_path),
                    "mime_type": PPTX_MIME_TYPE,
                    "slide_count": len(normalized_slides),
                },
                "metadata": metadata,
            }
        except Exception as exc:
            logger.error("create_presentation failed: %s", exc, exc_info=True)
            return self._error(str(exc))

    def edit_presentation_text(
        self,
        file_path: str,
        replacements: Any,
        output_filename: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Edit text in an existing PowerPoint using python-pptx.

        Args:
            file_path: Local path to an existing .pptx.
            replacements: Dict or JSON string mapping old text to new text.
            output_filename: Optional filename for the edited deck.
        """
        try:
            from pptx import Presentation
        except Exception:
            return self._error("python-pptx is not installed; add python-pptx to requirements.")

        try:
            source = Path(file_path)
            if not source.exists():
                return self._error(f"Presentation not found: {file_path}")

            mapping = _parse_jsonish(replacements, {})
            if not isinstance(mapping, dict) or not mapping:
                return self._error("replacements must be a non-empty dict or JSON object")

            prs = Presentation(str(source))
            changed = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            updated = run.text
                            for old, new in mapping.items():
                                if str(old) in updated:
                                    updated = updated.replace(str(old), str(new))
                            if updated != run.text:
                                run.text = updated
                                changed += 1

            safe_name = _safe_slug(output_filename or f"{source.stem}-edited")
            if not safe_name.endswith(".pptx"):
                safe_name = f"{safe_name}.pptx"
            output_path = source.parent / safe_name
            prs.save(str(output_path))

            artifact_id = self._persist_pptx(
                file_path=str(output_path),
                filename=safe_name,
                file_content=output_path.read_bytes(),
            )
            download_url = None
            if artifact_id and self.user_id:
                download_url = get_persistence_service().get_artifact_download_url(
                    artifact_id=artifact_id,
                    user_id=str(self.user_id),
                    expiry=3600,
                )

            return {
                "ok": True,
                "message": f"Edited {changed} text run(s) and saved '{safe_name}'.",
                "data": {
                    "artifact_id": artifact_id,
                    "filename": safe_name,
                    "download_url": download_url,
                    "local_path": str(output_path),
                    "changed_runs": changed,
                },
                "metadata": {
                    "kind": "presentation_tool_output",
                    "action": "edit_presentation_text",
                    "preview_type": "presentation",
                    "output_id": artifact_id or str(uuid.uuid4()),
                    "artifact_id": artifact_id,
                    "title": safe_name,
                    "summary": f"Edited {changed} text run(s) in a native PowerPoint file.",
                    "filename": safe_name,
                    "mime_type": PPTX_MIME_TYPE,
                    "download_url": download_url,
                    "inline": {
                        "topic": safe_name,
                        "slide_count": len(prs.slides),
                        "slides": [
                            {"index": i + 1, "title": self._first_slide_text(slide)}
                            for i, slide in enumerate(prs.slides)
                        ],
                    },
                },
            }
        except Exception as exc:
            logger.error("edit_presentation_text failed: %s", exc, exc_info=True)
            return self._error(str(exc))

    def _first_slide_text(self, slide: Any) -> str:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text:
                return str(shape.text).strip().splitlines()[0][:120]
        return "Slide"

    def _persist_pptx(self, *, file_path: str, filename: str, file_content: bytes) -> Optional[str]:
        if not (self.user_id and self.session_id):
            return None
        try:
            persistence = get_persistence_service()
            execution_id = persistence.create_execution_record(
                user_id=str(self.user_id),
                session_id=str(self.session_id),
                sandbox_id="presentation-tools",
                command=f"create_presentation {filename}",
                message_id=self.message_id,
            )
            if not execution_id:
                logger.warning("Could not create execution record for presentation artifact")
                return None

            artifact_id = persistence.create_artifact(
                execution_id=execution_id,
                user_id=str(self.user_id),
                session_id=str(self.session_id),
                sandbox_id="presentation-tools",
                file_path=file_path,
                file_content=file_content,
                mime_type=mimetypes.guess_type(filename)[0] or PPTX_MIME_TYPE,
                message_id=self.message_id,
            )
            try:
                persistence.db.table("sandbox_executions").update({
                    "status": "COMPLETED" if artifact_id else "FAILED",
                    "exit_code": 0 if artifact_id else 1,
                    "finished_at": datetime.utcnow().isoformat(),
                }).eq("execution_id", execution_id).execute()
            except Exception as exc:
                logger.warning("Failed to finalize presentation execution %s: %s", execution_id, exc)
            return artifact_id
        except Exception as exc:
            logger.warning("Failed to persist presentation artifact: %s", exc, exc_info=True)
            return None

    def _emit_presentation_created(self, metadata: Dict[str, Any]) -> None:
        if not (self.socketio and self.session_id):
            return
        payload = {
            "id": self.message_id,
            "conversationId": self.session_id,
            "metadata": metadata,
            "agent_name": "presentation_agent",
        }
        try:
            self.socketio.emit("presentation_generated", payload, room=f"conv:{self.session_id}")
        except Exception as exc:
            logger.warning("Failed to emit presentation_generated: %s", exc)

    def _error(self, message: str) -> Dict[str, Any]:
        return {
            "ok": False,
            "message": f"Presentation tool error: {message}",
            "error": message,
            "metadata": {
                "kind": "presentation_tool_output",
                "action": "error",
                "preview_type": "text",
                "title": "Presentation error",
                "summary": message,
                "inline": {"text_preview": message},
            },
        }


def build_presentation_agent(
    *,
    user_id: Optional[str],
    session_id: Optional[str],
    message_id: Optional[str],
    socketio=None,
    sid: Optional[str] = None,
    debug_mode: bool = True,
) -> Agent:
    tools = [
        PresentationTools(
            user_id=user_id,
            session_id=session_id,
            message_id=message_id,
            socketio=socketio,
            sid=sid,
        )
    ]
    from openrouter_reasoning_model import get_openrouter_model
    return Agent(
        name="presentation_agent",
        model=get_openrouter_model("xiaomi/mimo-v2.5"),
        role=(
            "Native PowerPoint specialist. Plans concise decks and creates editable "
            ".pptx files using presentation_tools."
        ),
        tools=tools,
        instructions=[
            "<system_instructions>",
            "You create native editable PowerPoint presentations, not HTML pages or image-only slide decks.",
            "If the user or Aetheria provides a hidden presentation template instruction, call create_presentation with that exact template id.",
            "Use list_presentation_templates only when template fit is unclear; it returns compact summaries to save context.",
            "Use get_presentation_template_details only for the one template you plan to use when you need its detailed design/layout guidance.",
            "For create_presentation, provide structured slides with types, titles, bullets, metrics, charts, tables, diagrams, visual summaries, and speaker notes where useful.",
            "Do not make a deck that is only title plus plain bullet slides. Use the backend template layouts: cover, insight cards, comparison, evidence chart, table, process/diagram, and visual explanation.",
            "For venture_blueprint, write like a premium business or pitch deck: title, problem/solution, market evidence, product/vision, business model, roadmap, and ask. Use short text blocks that fit the designed regions.",
            "When making comparison slides, always provide left/right titles and left/right bullet content. When making chart slides, provide chart.data. When making process slides, provide nodes or steps.",
            "Prefer concise claim-style titles and 3-6 strong slides unless the user asks for a different length.",
            "Use chart.data for simple bar evidence, table for comparison rows, nodes/steps for workflow diagrams, and metrics for rails.",
            "After create_presentation returns, inspect metadata.layout_validation when present. If it reports overflow, out-of-bounds, or overlap errors, regenerate with shorter titles/bullets or a better slide type before presenting the final answer.",
            "Return the artifact result naturally and mention that the file is downloadable and editable in PowerPoint.",
            "</system_instructions>",
        ],
        debug_mode=debug_mode,
    )
